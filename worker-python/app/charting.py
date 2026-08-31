import io
import math
import re
import zipfile
from pathlib import Path
from typing import Optional, Sequence

from PIL import Image, ImageDraw, ImageFont
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from .models import DeliverableChart


PALETTE = ["#2F6BCA", "#26A69A", "#F4B740", "#D95763", "#7656B6", "#607D8B"]
PPT_FONT = "Microsoft YaHei"


def valid_chart(chart: Optional[DeliverableChart]) -> bool:
    if chart is None or not chart.categories or not chart.series:
        return False
    category_count = len(chart.categories)
    return any(len(series.values) == category_count and series.values for series in chart.series)


def add_native_ppt_chart(slide, chart: DeliverableChart, left: float, top: float,
                         width: float, height: float, accent: str = "2F6BCA"):
    data = ChartData()
    data.categories = chart.categories
    usable = [item for item in chart.series if len(item.values) == len(chart.categories)]
    for series in usable:
        data.add_series(series.name or "指标", series.values)
    chart_type = {
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    }.get(chart.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    shape = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), data)
    native = shape.chart
    native.has_title = True
    native.chart_title.text_frame.text = chart.title
    title_font = native.chart_title.text_frame.paragraphs[0].font
    title_font.size = Pt(18)
    title_font.name = PPT_FONT
    native.has_legend = len(usable) > 1 or chart.type == "pie"
    if native.has_legend:
        native.legend.position = XL_LEGEND_POSITION.BOTTOM
        native.legend.include_in_layout = False
        native.legend.font.name = PPT_FONT
        native.legend.font.size = Pt(10)
    native.style = 10
    try:
        native.category_axis.tick_labels.font.name = PPT_FONT
        native.category_axis.tick_labels.font.size = Pt(10)
        native.value_axis.tick_labels.font.name = PPT_FONT
        native.value_axis.tick_labels.font.size = Pt(10)
    except (AttributeError, ValueError):
        pass
    for index, series in enumerate(native.series):
        color = RGBColor.from_string(accent if index == 0 else PALETTE[index % len(PALETTE)].lstrip("#"))
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.color.rgb = color
        except (AttributeError, ValueError):
            pass
    return shape


def chart_png(chart: DeliverableChart, width: int = 1200, height: int = 650) -> bytes:
    image = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(image)
    title_font = _font(30)
    label_font = _font(20)
    small_font = _font(17)
    draw.text((55, 28), chart.title, fill="#1F2937", font=title_font)
    plot = (85, 105, width - 65, height - 105)
    if chart.type == "pie":
        _draw_pie(draw, chart, plot, label_font, small_font)
    else:
        _draw_axes_chart(draw, chart, plot, label_font, small_font)
    stream = io.BytesIO()
    image.save(stream, format="PNG", optimize=True)
    return stream.getvalue()


def normalize_ppt_chart_ids(content: bytes) -> bytes:
    source = io.BytesIO(content)
    target = io.BytesIO()
    pattern = re.compile(rb'(<c:(?:axId|crossAx)\b[^>]*\bval=")(-\d+)(")')
    with zipfile.ZipFile(source, "r") as source_zip, zipfile.ZipFile(target, "w") as target_zip:
        for item in source_zip.infolist():
            data = source_zip.read(item.filename)
            if item.filename.startswith("ppt/charts/") and item.filename.endswith(".xml"):
                data = pattern.sub(lambda match: match.group(1) + str(int(match.group(2)) + 2 ** 32).encode("ascii")
                                   + match.group(3), data)
            target_zip.writestr(item, data)
    return target.getvalue()


def _draw_axes_chart(draw: ImageDraw.ImageDraw, chart: DeliverableChart, plot: tuple[int, int, int, int],
                     label_font, small_font) -> None:
    left, top, right, bottom = plot
    usable = [item for item in chart.series if len(item.values) == len(chart.categories)][:4]
    values = [value for item in usable for value in item.values]
    low = min(0.0, min(values, default=0.0))
    high = max(values, default=1.0)
    if math.isclose(low, high):
        high = low + 1.0
    draw.line((left, top, left, bottom), fill="#A9B7C6", width=2)
    draw.line((left, bottom, right, bottom), fill="#A9B7C6", width=2)
    for step in range(5):
        y = bottom - (bottom - top) * step / 4
        value = low + (high - low) * step / 4
        draw.line((left, y, right, y), fill="#E5ECF3", width=1)
        draw.text((8, y - 10), _number(value), fill="#718398", font=small_font)
    count = max(1, len(chart.categories))
    slot = (right - left) / count
    colors = PALETTE
    if chart.type == "line":
        for series_index, series in enumerate(usable):
            points = []
            for index, value in enumerate(series.values):
                x = left + slot * (index + 0.5)
                y = bottom - (value - low) / (high - low) * (bottom - top)
                points.append((x, y))
            if len(points) > 1:
                draw.line(points, fill=colors[series_index], width=5, joint="curve")
            for x, y in points:
                draw.ellipse((x - 6, y - 6, x + 6, y + 6), fill=colors[series_index])
    else:
        bar_width = min(54, slot * 0.72 / max(1, len(usable)))
        for category_index in range(count):
            for series_index, series in enumerate(usable):
                value = series.values[category_index]
                zero_y = bottom - (0 - low) / (high - low) * (bottom - top)
                value_y = bottom - (value - low) / (high - low) * (bottom - top)
                center = left + slot * (category_index + 0.5)
                x = center + (series_index - (len(usable) - 1) / 2) * bar_width
                draw.rounded_rectangle((x - bar_width * .42, min(zero_y, value_y), x + bar_width * .42,
                                        max(zero_y, value_y)), radius=4, fill=colors[series_index])
    for index, category in enumerate(chart.categories):
        x = left + slot * (index + 0.5)
        label = str(category)[:10]
        box = draw.textbbox((0, 0), label, font=small_font)
        draw.text((x - (box[2] - box[0]) / 2, bottom + 14), label, fill="#52677E", font=small_font)
    _draw_legend(draw, usable, right, top, label_font)


def _draw_pie(draw: ImageDraw.ImageDraw, chart: DeliverableChart, plot: tuple[int, int, int, int],
              label_font, small_font) -> None:
    series = next((item for item in chart.series if len(item.values) == len(chart.categories)), None)
    if not series:
        return
    values = [max(0.0, value) for value in series.values]
    total = sum(values) or 1.0
    left, top, right, bottom = plot
    size = min(bottom - top, (right - left) * .58)
    box = (left + 50, top + 12, left + 50 + size, top + 12 + size)
    angle = -90.0
    for index, value in enumerate(values):
        extent = 360.0 * value / total
        draw.pieslice(box, angle, angle + extent, fill=PALETTE[index % len(PALETTE)], outline="white", width=3)
        angle += extent
    legend_x = int(left + size + 105)
    for index, (category, value) in enumerate(zip(chart.categories, values)):
        y = top + 35 + index * 48
        draw.rounded_rectangle((legend_x, y, legend_x + 24, y + 24), radius=4,
                               fill=PALETTE[index % len(PALETTE)])
        draw.text((legend_x + 38, y - 2), f"{str(category)[:16]}  {value / total:.1%}",
                  fill="#344A60", font=label_font)


def _draw_legend(draw: ImageDraw.ImageDraw, series: Sequence, right: int, top: int, font) -> None:
    if len(series) <= 1:
        return
    x = right - 220
    for index, item in enumerate(series):
        y = top + index * 34
        draw.rectangle((x, y, x + 20, y + 20), fill=PALETTE[index % len(PALETTE)])
        draw.text((x + 30, y - 3), item.name[:12], fill="#52677E", font=font)


def _font(size: int):
    for path in (
        Path("/System/Library/Fonts/PingFang.ttc"),
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    ):
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def _number(value: float) -> str:
    absolute = abs(value)
    if absolute >= 100_000_000:
        return f"{value / 100_000_000:.1f}亿"
    if absolute >= 10_000:
        return f"{value / 10_000:.1f}万"
    return f"{value:,.1f}".rstrip("0").rstrip(".")
