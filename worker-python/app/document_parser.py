import csv
import hashlib
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .config import settings
from .models import ParsedChunk, ParsedDocument


@dataclass
class LocatedText:
    text: str
    location: Dict[str, object]


def parse_document(path: Path, file_name: str, media_type: str) -> ParsedDocument:
    suffix = Path(file_name).suffix.lower()
    warnings: List[str] = []
    if suffix == ".pdf" or media_type == "application/pdf":
        located = _parse_pdf(path)
    elif suffix == ".docx":
        located = _parse_docx(path)
    elif suffix == ".pptx":
        located = _parse_pptx(path)
    elif suffix in {".xlsx", ".xlsm"}:
        located, warnings = _parse_xlsx(path)
    elif suffix in {".csv", ".tsv"}:
        located, warnings = _parse_delimited(path, suffix == ".tsv")
    elif suffix in {".txt", ".md", ".mermaid", ".mmd", ".excalidraw", ".srt", ".vtt", ".json"} or media_type.startswith("text/"):
        located = [LocatedText(_read_text(path), {"type": "text", "start": 0})]
    elif suffix in {".mp4", ".mov", ".mkv", ".mp3", ".wav", ".m4a"}:
        located = _transcribe_media(path)
    else:
        raise ValueError("暂不支持该文件格式")

    chunks = _chunk_located_text(located)
    return ParsedDocument(
        file_name=file_name,
        media_type=media_type or "application/octet-stream",
        title=Path(file_name).stem,
        text_length=sum(len(item.text) for item in located),
        chunks=chunks,
        warnings=warnings,
    )


def _parse_pdf(path: Path) -> List[LocatedText]:
    reader = PdfReader(str(path), strict=False)
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exception:
            raise ValueError("PDF 已加密，无法解析") from exception
    result: List[LocatedText] = []
    for index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            result.append(LocatedText(text, {"type": "pdf", "page": index + 1}))
    return result


def _parse_docx(path: Path) -> List[LocatedText]:
    document = Document(str(path))
    result: List[LocatedText] = []
    for index, paragraph in enumerate(document.paragraphs):
        if paragraph.text.strip():
            result.append(LocatedText(paragraph.text, {"type": "word", "paragraph": index + 1}))
    for table_index, table in enumerate(document.tables):
        rows = ["\t".join(cell.text for cell in row.cells) for row in table.rows]
        if rows:
            result.append(LocatedText("\n".join(rows), {"type": "word_table", "table": table_index + 1}))
    return result


def _parse_pptx(path: Path) -> List[LocatedText]:
    presentation = Presentation(str(path))
    result: List[LocatedText] = []
    for slide_index, slide in enumerate(presentation.slides):
        text_parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)
        if text_parts:
            result.append(LocatedText("\n".join(text_parts), {"type": "powerpoint", "slide": slide_index + 1}))
    return result


def _parse_xlsx(path: Path) -> Tuple[List[LocatedText], List[str]]:
    workbook = load_workbook(str(path), read_only=True, data_only=False, keep_links=False)
    result: List[LocatedText] = []
    warnings: List[str] = []
    max_rows = 10_000
    for sheet in workbook.worksheets:
        rows: List[str] = []
        for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if row_index > max_rows:
                warnings.append("工作表 %s 超过 %d 行，作为知识资料仅解析前 %d 行" % (sheet.title, max_rows, max_rows))
                break
            rows.append("\t".join("" if value is None else str(value) for value in row))
        if rows:
            result.append(LocatedText("\n".join(rows), {"type": "excel", "sheet": sheet.title, "rows": "1:%d" % len(rows)}))
    workbook.close()
    return result, warnings


def _parse_delimited(path: Path, tsv: bool) -> Tuple[List[LocatedText], List[str]]:
    delimiter = "\t" if tsv else ","
    rows: List[str] = []
    warnings: List[str] = []
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as stream:
        reader = csv.reader(stream, delimiter=delimiter)
        for index, row in enumerate(reader, start=1):
            if index > 10_000:
                warnings.append("表格超过 10000 行，作为知识资料仅解析前 10000 行")
                break
            rows.append("\t".join(row))
    return [LocatedText("\n".join(rows), {"type": "table", "rows": "1:%d" % len(rows)})], warnings


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8-sig", errors="replace")


def _transcribe_media(path: Path) -> List[LocatedText]:
    try:
        from faster_whisper import WhisperModel
    except ImportError as exception:
        raise ValueError("当前部署未安装音视频转写组件 faster-whisper") from exception
    try:
        model = WhisperModel(settings.whisper_model, device=settings.whisper_device,
                             compute_type=settings.whisper_compute_type)
        segments, _ = model.transcribe(str(path), vad_filter=True, beam_size=5)
        result = []
        for segment in segments:
            text = segment.text.strip()
            if text:
                result.append(LocatedText(text, {
                    "type": "media",
                    "start_seconds": round(segment.start, 3),
                    "end_seconds": round(segment.end, 3),
                }))
        if not result:
            raise ValueError("音视频中没有识别到可转写内容")
        return result
    except ValueError:
        raise
    except Exception as exception:
        raise ValueError("音视频转写失败，请检查模型文件和媒体编码") from exception


def _chunk_located_text(items: Iterable[LocatedText]) -> List[ParsedChunk]:
    result: List[ParsedChunk] = []
    for item in items:
        normalized = re.sub(r"[ \t]+", " ", item.text).strip()
        if not normalized:
            continue
        start = 0
        while start < len(normalized):
            end = min(len(normalized), start + settings.knowledge_chunk_chars)
            if end < len(normalized):
                boundary = max(normalized.rfind("\n", start, end), normalized.rfind("。", start, end))
                if boundary > start + settings.knowledge_chunk_chars // 2:
                    end = boundary + 1
            text = normalized[start:end].strip()
            if text:
                location = dict(item.location)
                location["char_start"] = start
                location["char_end"] = end
                result.append(
                    ParsedChunk(
                        index=len(result),
                        text=text,
                        location=location,
                        content_hash="sha256:" + hashlib.sha256(text.encode("utf-8")).hexdigest(),
                    )
                )
            if end >= len(normalized):
                break
            start = max(start + 1, end - settings.knowledge_chunk_overlap)
    return result
