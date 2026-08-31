import io
import math
import re
from typing import Iterable, List, Optional

from lxml.etree import SubElement
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from .charting import add_native_ppt_chart, normalize_ppt_chart_ids, valid_chart
from .citations import inline_sources, reference_entries
from .models import DeliverableRequest, DeliverableSection


HUAWEI_STYLE_C = "guizang-huawei-style-c"
FRONTEND_SLIDES = "frontend-slides"
RED = RGBColor(199, 0, 11)
RED_DARK = RGBColor(153, 0, 8)
RED_SOFT = RGBColor(245, 229, 231)
PAPER = RGBColor(247, 248, 250)
CARD = RGBColor(255, 255, 255)
INK = RGBColor(34, 37, 43)
MUTED = RGBColor(111, 118, 128)
LINE = RGBColor(223, 227, 232)
GREY = RGBColor(240, 242, 245)
WHITE = RGBColor(255, 255, 255)
FONT_FAMILY = "Microsoft YaHei"
SLIDE_LEFT = 0.95
SLIDE_RIGHT = 12.38
CONTENT_TOP = 1.68
CONTENT_BOTTOM = 6.56


def catalog() -> list[dict[str, object]]:
    return [{
        "id": HUAWEI_STYLE_C,
        "name": "华为企业汇报 · Style C",
        "description": "白灰底、企业红、清晰层级，适合管理层汇报、经营复盘和工作计划。",
        "formats": ["pptx"],
        "theme": "Classic Red",
        "source": "https://github.com/SeanDongX/guizang-ppt-skill",
    }, {
        "id": FRONTEND_SLIDES,
        "name": "Frontend Slides · 网页演示",
        "description": "单文件 HTML + JavaScript，支持动效、键盘翻页和浏览器内播放；不是 PowerPoint 文件。",
        "formats": ["html_slides"],
        "theme": "Editorial Data",
        "source": "https://github.com/zarazhangrui/frontend-slides",
    }]


def render(request: DeliverableRequest) -> bytes:
    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    _cover(deck, request)
    slide_number = 2
    for section_index, section in enumerate(request.sections, start=1):
        points = _content_points(section)
        groups = [points[:3]] if valid_chart(section.chart) else list(_chunks(points, 4))
        for page_index, group in enumerate(groups, start=1):
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            _background(slide, PAPER)
            heading = section.heading if page_index == 1 else f"{section.heading}（续）"
            _chrome(slide, slide_number, section_index)
            _text(slide, heading, 1.82, 0.62, 10.45, 0.58, 27, INK, bold=True,
                  min_size=22, single_line=True)
            if valid_chart(section.chart):
                _chart_layout(slide, section.chart, group)
            else:
                _render_content(slide, heading, group, section_index)
            refs = inline_sources(request, section.refs)
            if refs and not (valid_chart(section.chart) and section.chart.source_ref):
                _text(slide, "参考：" + refs[:180], 0.95, 6.93, 10.55, 0.2, 8, MUTED,
                      min_size=8, single_line=True)
            slide_number += 1
    _reference_slides(deck, request, slide_number)
    stream = io.BytesIO()
    deck.save(stream)
    return normalize_ppt_chart_ids(stream.getvalue())


def _cover(deck: Presentation, request: DeliverableRequest) -> None:
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    _background(slide, PAPER)
    _rect(slide, 0, 0, 0.34, 7.5, RED)
    _rect(slide, 0.92, 1.18, 0.18, 0.18, RED)
    _text(slide, "MANAGEMENT REPORT", 1.28, 1.12, 4.5, 0.3, 10, RED, bold=True)
    _text(slide, request.title, 0.95, 2.03, 10.75, 1.48, 42, INK, bold=True, min_size=32)
    _text(slide, request.subtitle or "由 FinFlow Studio 工作流生成", 0.98, 3.7, 10.2, 0.5, 17, MUTED,
          min_size=14, single_line=True)
    _rect(slide, 0.98, 4.46, 1.36, 0.07, RED)
    _text(slide, "FINFLOW STUDIO", 10.75, 0.58, 1.65, 0.25, 10, INK, bold=True, align=PP_ALIGN.RIGHT)
    _text(slide, "企业汇报 · Classic Red", 0.98, 6.72, 4.2, 0.25, 10, MUTED)
    _text(slide, "01", 11.65, 5.64, 0.65, 0.5, 18, WHITE, bold=True, align=PP_ALIGN.CENTER,
          fill=RED, vertical=MSO_ANCHOR.MIDDLE)


def _chrome(slide, page_number: int, section_number: int) -> None:
    _rect(slide, 0.95, 0.78, 0.18, 0.18, RED)
    _text(slide, f"{section_number:02d}", 1.22, 0.7, 0.45, 0.3, 10, RED, bold=True,
          min_size=10, single_line=True)
    _text(slide, "FINFLOW STUDIO", 10.72, 0.26, 1.65, 0.22, 9, INK, bold=True,
          align=PP_ALIGN.RIGHT, min_size=9, single_line=True)
    _rect(slide, 0.95, 1.38, 11.42, 0.018, LINE)
    _text(slide, str(page_number), 11.98, 6.96, 0.35, 0.18, 8, MUTED,
          align=PP_ALIGN.RIGHT, min_size=8, single_line=True)


def _cards(slide, points: List[str]) -> None:
    positions = [(0.95, 1.8), (6.72, 1.8), (0.95, 4.28), (6.72, 4.28)]
    for index, point in enumerate(points):
        left, top = positions[index]
        _rect(slide, left, top, 5.25, 2.05, CARD, line=LINE)
        _rect(slide, left, top, 0.08, 2.05, RED)
        _text(slide, f"{index + 1:02d}", left + 0.36, top + 0.28, 0.56, 0.34, 13, RED, bold=True)
        _text(slide, point, left + 0.36, top + 0.79, 4.45, 0.98, 16, INK, bold=False)


def _render_content(slide, heading: str, points: List[str], section_index: int) -> None:
    normalized = heading.lower()
    if section_index == 1 or any(keyword in normalized for keyword in ("摘要", "总体", "核心结论")):
        _summary_layout(slide, points)
    elif any(keyword in normalized for keyword in ("行动", "责任", "闭环", "计划", "要求")):
        _action_layout(slide, points)
    elif any(keyword in normalized for keyword in ("top", "物料", "清单", "事项")):
        _ranked_layout(slide, points)
    elif any(keyword in normalized for keyword in ("仓", "品类", "成因", "原因")):
        _split_layout(slide, points)
    elif any(keyword in normalized for keyword in ("差异", "库龄", "减值", "金额", "占比")):
        _metric_layout(slide, points)
    else:
        _split_layout(slide, points)


def _summary_layout(slide, points: List[str]) -> None:
    lead = points[0] if points else "工作流已完成本次分析。"
    _rect(slide, 0.95, 1.76, 0.1, 4.72, RED)
    _text(slide, "核心判断", 1.34, 1.9, 1.25, 0.3, 11, RED, bold=True,
          min_size=11, single_line=True)
    _text(slide, lead, 1.32, 2.48, 5.6, 2.2, 25, INK, bold=True, min_size=19,
          vertical=MSO_ANCHOR.MIDDLE)
    _rect(slide, 1.34, 5.24, 1.34, 0.07, RED)
    supporting = points[1:4] or [lead]
    for index, point in enumerate(supporting):
        top = 1.83 + index * 1.52
        _text(slide, f"0{index + 1}", 7.38, top + 0.04, 0.42, 0.28, 11, RED, bold=True,
              min_size=11, single_line=True)
        _text(slide, point, 8.05, top, 4.0, 0.92, 16, INK, min_size=14)
        if index < len(supporting) - 1:
            _rect(slide, 7.38, top + 1.12, 4.68, 0.015, LINE)


def _metric_layout(slide, points: List[str]) -> None:
    lead = points[0] if points else "暂无指标"
    metric = _first_metric(lead)
    _rect(slide, 0.95, 1.78, 4.05, 4.82, CARD, line=LINE)
    _rect(slide, 0.95, 1.78, 0.1, 4.82, RED)
    _text(slide, "关键指标", 1.38, 2.17, 1.5, 0.28, 11, MUTED, bold=True,
          min_size=11, single_line=True)
    _text(slide, metric or "重点关注", 1.35, 2.78, 3.2, 0.72, 30 if metric else 23, RED_DARK,
          bold=True, min_size=22, single_line=True)
    lead_text = lead.replace(metric, "", 1).strip("，,：:；; ") if metric else lead
    _text(slide, lead_text or lead, 1.38, 3.78, 3.12, 1.55, 17, INK, bold=True, min_size=15)
    _text(slide, "数据口径以工作流上游结果为准", 1.38, 5.85, 3.1, 0.28, 10, MUTED)
    supporting = points[1:4]
    for index, point in enumerate(supporting):
        top = 1.88 + index * 1.53
        _text(slide, f"0{index + 1}", 5.55, top + 0.08, 0.55, 0.3, 12, RED, bold=True)
        _text(slide, point, 6.25, top, 5.72, 0.92, 16, INK, min_size=14)
        if index < len(supporting) - 1:
            _rect(slide, 5.55, top + 1.16, 6.42, 0.015, LINE)


def _ranked_layout(slide, points: List[str]) -> None:
    for index, point in enumerate(points[:4]):
        top = 1.78 + index * 1.2
        fill = CARD if index % 2 == 0 else GREY
        _rect(slide, 0.95, top, 11.42, 0.98, fill, line=LINE if index % 2 == 0 else None)
        _text(slide, str(index + 1), 1.18, top + 0.17, 0.56, 0.56, 18, WHITE, bold=True,
              align=PP_ALIGN.CENTER, fill=RED if index == 0 else RED_DARK, vertical=MSO_ANCHOR.MIDDLE)
        _text(slide, point, 2.02, top + 0.16, 9.75, 0.62, 17, INK, bold=index == 0,
              min_size=15,
              vertical=MSO_ANCHOR.MIDDLE)


def _action_layout(slide, points: List[str]) -> None:
    items = points[:4]
    if not items:
        items = ["明确责任人与完成时间"]
    _rect(slide, 1.25, 3.08, 10.5, 0.04, LINE)
    width = 10.9 / len(items)
    for index, point in enumerate(items):
        left = 0.95 + index * width
        _text(slide, f"0{index + 1}", left + 0.18, 2.57, 0.62, 0.62, 16, WHITE, bold=True,
              align=PP_ALIGN.CENTER, fill=RED if index == 0 else RED_DARK, vertical=MSO_ANCHOR.MIDDLE)
        _text(slide, "责任动作", left + 0.18, 3.48, width - 0.36, 0.3, 11, RED, bold=True)
        _text(slide, point, left + 0.18, 4.0, width - 0.38, 1.7, 16, INK, bold=index == 0,
              min_size=14)


def _split_layout(slide, points: List[str]) -> None:
    lead = points[0] if points else "工作流已完成本次分析。"
    _rect(slide, 0.95, 1.78, 4.62, 4.82, INK)
    _text(slide, "业务判断", 1.32, 2.18, 1.3, 0.28, 11, RED_SOFT, bold=True)
    _text(slide, lead, 1.32, 2.82, 3.85, 2.15, 23, WHITE, bold=True, min_size=18)
    _rect(slide, 1.32, 5.58, 1.15, 0.06, RED)
    for index, point in enumerate(points[1:4]):
        top = 1.93 + index * 1.48
        _text(slide, f"0{index + 1}", 6.12, top, 0.55, 0.3, 12, RED, bold=True)
        _text(slide, point, 6.82, top - 0.02, 5.0, 0.92, 16, INK, min_size=14)
        if index < min(2, len(points[1:4]) - 1):
            _rect(slide, 6.12, top + 1.12, 5.7, 0.015, LINE)


def _chart_layout(slide, chart, points: List[str]) -> None:
    add_native_ppt_chart(slide, chart, 0.95, 1.72, 7.35, 4.82, accent="C7000B")
    _text(slide, "关键判断", 8.72, 1.82, 2.2, 0.3, 11, RED, bold=True)
    for index, point in enumerate(points[:3]):
        top = 2.35 + index * 1.34
        _text(slide, f"0{index + 1}", 8.72, top + 0.02, 0.45, 0.3, 11, RED, bold=True)
        _text(slide, point, 9.28, top, 2.92, 1.02, 15, INK, bold=index == 0, min_size=13)
    if chart.source_ref:
        _text(slide, "图表来源：" + chart.source_ref, 0.98, 6.78, 7.2, 0.2, 8, MUTED,
              min_size=8, single_line=True)


def _first_metric(value: str) -> str:
    match = re.search(r"[-+]?\d[\d,.]*(?:%|亿元|万元|万|天)", value)
    return match.group(0) if match else ""


def _background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor,
          line: Optional[RGBColor] = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def _text(slide, text: str, left: float, top: float, width: float, height: float, size: int,
          color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT, fill: Optional[RGBColor] = None,
          vertical=MSO_ANCHOR.TOP, min_size: int = 12, single_line: bool = False) -> None:
    box = _rect(slide, left, top, width, height, fill) if fill else slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    frame.margin_left = frame.margin_right = Inches(0.05) if fill else 0
    frame.margin_top = frame.margin_bottom = 0
    normalized = re.sub(r"\s*\n+\s*", " " if single_line else "\n", str(text)).strip()
    fitted_size = _fit_font_size(normalized, width, height, size, min_size, single_line)
    for line_index, line in enumerate(normalized.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = 1.08
        paragraph.space_after = Pt(0)
        parts = re.split(r"(\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\))", line)
        for part in parts:
            if not part:
                continue
            run = paragraph.add_run()
            run.text = part
            citation = bool(re.fullmatch(r"\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\)", part))
            run.font.size = Pt(max(7, round(fitted_size * 0.58))) if citation else Pt(fitted_size)
            run.font.bold = bold
            run.font.color.rgb = color
            _set_run_font(run, FONT_FAMILY)


def _fit_font_size(text: str, width: float, height: float, preferred: int, minimum: int,
                   single_line: bool) -> int:
    visual_units = sum(1.0 if ord(char) > 255 else 0.56 for char in re.sub(r"\s+", " ", text))
    explicit_lines = max(1, text.count("\n") + 1)
    for candidate in range(preferred, minimum - 1, -1):
        units_per_line = max(1.0, width * 72 / (candidate * 0.94))
        estimated_lines = explicit_lines if single_line else max(explicit_lines, math.ceil(visual_units / units_per_line))
        if single_line and visual_units > units_per_line:
            continue
        if estimated_lines * candidate * 1.18 <= height * 72:
            return candidate
    return minimum


def _set_run_font(run, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = properties.find(qn(tag))
        if element is None:
            element = SubElement(properties, qn(tag))
        element.set("typeface", family)


def _reference_slides(deck: Presentation, request: DeliverableRequest, page_number: int) -> None:
    entries = reference_entries(request)
    for page_index, group in enumerate(_chunks(entries, 8), start=1):
        slide = deck.slides.add_slide(deck.slide_layouts[6])
        _background(slide, PAPER)
        _chrome(slide, page_number, 99)
        heading = "参考文献" if page_index == 1 else "参考文献（续）"
        _text(slide, heading, 1.82, 0.62, 10.45, 0.58, 27, INK, bold=True,
              min_size=22, single_line=True)
        for index, entry in enumerate(group):
            _text(slide, entry, 1.1, 1.72 + index * 0.62, 10.9, 0.48, 11, INK, min_size=9)
        page_number += 1


def _content_points(section: DeliverableSection) -> List[str]:
    raw: List[str] = []
    for paragraph in section.paragraphs:
        raw.extend(re.split(r"\n+|(?<=[。！？；])", paragraph))
    raw.extend(section.bullets)
    result: List[str] = []
    for value in raw:
        clean = value.replace("**", "").replace("`", "")
        clean = re.sub(r"^\s*(?:[#>*•\-]+\s*|\d+[.、]\s+)+", "", clean)
        clean = re.sub(r"\s+", " ", clean).strip(" ；;。")
        if len(clean) >= 4 and clean not in result:
            result.extend(item for item in _shorten_point(clean) if item not in result)
        if len(result) >= 36:
            break
    return result or ["工作流已完成处理，但当前步骤没有生成可展示的文字内容。"]


def _shorten_point(value: str, limit: int = 72) -> List[str]:
    if len(value) <= limit:
        return [value]
    clauses = [item.strip() for item in re.split(r"(?<=[，,；;。])", value) if item.strip()]
    result: List[str] = []
    current = ""
    for clause in clauses:
        if current and len(current) + len(clause) > limit:
            result.append(current.rstrip("，,；;。"))
            current = clause
        else:
            current += clause
    if current:
        result.append(current.rstrip("，,；;。"))
    return [item[:limit] for item in result[:2]]


def _chunks(values: List[str], size: int) -> Iterable[List[str]]:
    for index in range(0, len(values), size):
        yield values[index:index + size]


def _format_refs(section: DeliverableSection) -> str:
    return "；".join(ref.source_name for ref in section.refs)
