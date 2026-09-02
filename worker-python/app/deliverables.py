import io
import json
import math
import re
from html import escape
from pathlib import Path
from typing import Iterable, List, Optional

from docx import Document
from lxml.etree import SubElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn as pptx_qn
from pptx.util import Inches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import Image as ReportLabImage, ListFlowable, ListItem, PageBreak, Paragraph, SimpleDocTemplate, Spacer

from .charting import add_native_ppt_chart, chart_png, normalize_ppt_chart_ids, valid_chart
from .citations import inline_sources, normalize_markers, reference_entries, reference_records
from .models import DeliverableChart, DeliverableRequest, DeliverableSection
from .ppt_skills import HUAWEI_STYLE_C, render as render_ppt_skill


BLUE = PptColor(47, 107, 202)
DARK_BLUE = PptColor(25, 72, 126)
DARK = PptColor(31, 41, 55)
MUTED = PptColor(91, 105, 125)
PALE_BLUE = PptColor(238, 246, 255)
LIGHT_LINE = PptColor(218, 229, 242)
PPT_FONT = "Microsoft YaHei"


def create_pptx(request: DeliverableRequest) -> bytes:
    request = _normalize_ppt_request(request)
    if request.ppt_skill == HUAWEI_STYLE_C:
        return render_ppt_skill(request)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_text(title_slide, request.title, 0.85, 2.15, 11.5, 1.15, 42, DARK, True, min_size=30)
    _add_text(title_slide, request.subtitle or "由 FinFlow Studio 工作流生成", 0.88, 3.45, 11.2, 0.55,
              20, MUTED, False, min_size=15, single_line=True)
    _add_brand(title_slide, 1)
    slide_number = 2
    for section in request.sections:
        points = _content_points(section)
        if valid_chart(section.chart):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            _add_slide_heading(slide, section.heading, slide_number)
            add_native_ppt_chart(slide, section.chart, 0.85, 1.72, 7.5, 4.7)
            _add_compact_points(slide, points[:3], 8.75, 1.9, 3.65)
            _add_brand(slide, slide_number)
            reference = inline_sources(request, section.refs)
            if reference:
                _add_text(slide, "参考：" + reference[:240], 0.85, 6.82, 11.45, 0.24, 9, MUTED, False)
            slide_number += 1
            continue
        for page_index, group in enumerate(_chunks(points, 4), start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            heading = section.heading if page_index == 1 else f"{section.heading}（续）"
            _add_slide_heading(slide, heading, slide_number)
            _add_point_list(slide, group)
            _add_brand(slide, slide_number)
            reference = inline_sources(request, section.refs)
            if reference:
                _add_text(slide, "参考：" + reference[:240], 0.85, 6.82, 11.45, 0.24, 9, MUTED, False)
            slide_number += 1
    _add_reference_slides(presentation, request, slide_number)
    stream = io.BytesIO()
    presentation.save(stream)
    return normalize_ppt_chart_ids(stream.getvalue())


def create_html_slides(request: DeliverableRequest) -> bytes:
    request = _normalize_ppt_request(request)
    slides: List[str] = [f"""
      <section class="slide cover active" data-kind="cover">
        <div class="cover-kicker">FINFLOW WEB PRESENTATION</div>
        <div class="cover-rule"></div>
        <h1>{escape(request.title)}</h1>
        <p>{escape(request.subtitle or '由 FinFlow Studio 工作流生成')}</p>
        <div class="cover-meta"><span>HTML + JavaScript</span><span>浏览器演示文件</span></div>
      </section>"""]
    charts: List[dict[str, object]] = []
    slide_number = 2
    for section_index, section in enumerate(request.sections, start=1):
        points = _content_points(section)
        groups = [points[:3]] if valid_chart(section.chart) else list(_chunks(points, 4))
        for page_index, group in enumerate(groups, start=1):
            heading = section.heading if page_index == 1 else f"{section.heading}（续）"
            source = inline_sources(request, section.refs)
            chart_index: Optional[int] = None
            if valid_chart(section.chart):
                chart_index = len(charts)
                charts.append(section.chart.model_dump())
            body = _html_slide_body(heading, group, section_index, chart_index)
            source_html = f'<div class="slide-source">参考：{escape(source[:180])}</div>' if source else ""
            slides.append(f"""
      <section class="slide" data-kind="content">
        <header class="slide-header"><span>{section_index:02d}</span><h2>{escape(heading)}</h2><em>FINFLOW STUDIO</em></header>
        {body}
        {source_html}<div class="slide-number">{slide_number:02d}</div>
      </section>""")
            slide_number += 1
    entries = reference_entries(request)
    for page_index, group in enumerate(_chunks(entries, 7), start=1):
        heading = "参考文献" if page_index == 1 else "参考文献（续）"
        items = "".join(f"<li>{escape(entry)}</li>" for entry in group)
        slides.append(f"""
      <section class="slide references" data-kind="references">
        <header class="slide-header"><span>REF</span><h2>{heading}</h2><em>FINFLOW STUDIO</em></header>
        <ol>{items}</ol><div class="slide-number">{slide_number:02d}</div>
      </section>""")
        slide_number += 1
    chart_data = json.dumps(charts, ensure_ascii=False, separators=(",", ":")).replace("</", "<\\/")
    html = _html_slides_document(request.title, "".join(slides), chart_data)
    return html.encode("utf-8")


def _html_slide_body(heading: str, points: List[str], section_index: int,
                     chart_index: Optional[int]) -> str:
    if chart_index is not None:
        insights = "".join(
            f'<li><span>{index:02d}</span><p>{_html_rich_text(point)}</p></li>'
            for index, point in enumerate(points[:3], start=1)
        )
        return f'<div class="layout chart-layout"><div class="chart" data-chart="{chart_index}"></div><div class="insights"><h3>关键判断</h3><ol>{insights}</ol></div></div>'
    lead = points[0] if points else "工作流已完成处理。"
    supporting = points[1:4]
    normalized = heading.lower()
    if section_index == 1 or any(keyword in normalized for keyword in ("摘要", "总体", "核心结论")):
        evidence = "".join(
            f'<li><span>{index:02d}</span><p>{_html_rich_text(point)}</p></li>'
            for index, point in enumerate(supporting, start=1)
        )
        return f'<div class="layout statement-layout"><div class="statement"><small>核心判断</small><h3>{_html_rich_text(lead)}</h3></div><ol class="evidence">{evidence}</ol></div>'
    if any(keyword in normalized for keyword in ("行动", "责任", "计划", "建议", "下一步")):
        items = "".join(
            f'<li><span>{index:02d}</span><p>{_html_rich_text(point)}</p></li>'
            for index, point in enumerate(points[:4], start=1)
        )
        return f'<div class="layout action-layout"><div class="action-line"></div><ol>{items}</ol></div>'
    metric = _first_metric(lead)
    if metric:
        detail = lead.replace(metric, "", 1).strip("，,：:；; ") or lead
        items = "".join(f'<li>{_html_rich_text(point)}</li>' for point in supporting)
        return f'<div class="layout metric-layout"><div class="metric"><small>关键指标</small><strong>{escape(metric)}</strong><p>{_html_rich_text(detail)}</p></div><ul>{items}</ul></div>'
    items = "".join(
        f'<li><span>{index:02d}</span><p>{_html_rich_text(point)}</p></li>'
        for index, point in enumerate(supporting, start=1)
    )
    return f'<div class="layout editorial-layout"><div class="editorial-lead"><small>业务判断</small><h3>{_html_rich_text(lead)}</h3></div><ol>{items}</ol></div>'


def _html_rich_text(value: str) -> str:
    parts = re.split(r"(\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\))", value)
    return "".join(
        f'<sup>{escape(part)}</sup>' if re.fullmatch(r"\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\)", part)
        else escape(part) for part in parts if part
    )


def _html_slides_document(title: str, slides: str, chart_data: str) -> str:
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<meta name="generator" content="FinFlow Studio · Frontend Slides inspired"><title>{escape(title)}</title>
<style>
:root{{--paper:#f7f8fa;--ink:#18202c;--muted:#647286;--line:#dce2e9;--blue:#1769e0;--teal:#0b9f91;--coral:#f05a47;--yellow:#f2bd42}}
*{{box-sizing:border-box}}html,body{{width:100%;height:100%;margin:0;overflow:hidden;background:#0f1722;color:var(--ink);font-family:Inter,"Microsoft YaHei","PingFang SC",Arial,sans-serif}}
.deck{{position:relative;width:100%;height:100%;display:grid;place-items:center}}.stage{{position:relative;width:min(100vw,calc(100vh * 16 / 9));height:min(100vh,calc(100vw * 9 / 16));overflow:hidden;background:var(--paper);box-shadow:0 28px 90px #0008}}
.slide{{position:absolute;inset:0;display:none;padding:5.8% 6.6% 5.2%;background:var(--paper);overflow:hidden}}.slide.active{{display:block}}.slide.active .layout>*{{animation:rise .55s cubic-bezier(.2,.8,.2,1) both}}.slide.active .layout>*:nth-child(2){{animation-delay:.09s}}@keyframes rise{{from{{opacity:0;transform:translateY(18px)}}to{{opacity:1;transform:none}}}}
.cover{{padding:8% 8.4%;background:linear-gradient(120deg,#f8fbff 0 68%,#e8f1ff 68%)}}.cover:after{{content:"";position:absolute;right:6%;bottom:7%;width:28%;height:9px;background:linear-gradient(90deg,var(--blue) 0 58%,var(--coral) 58% 78%,var(--yellow) 78%)}}.cover-kicker{{font-size:clamp(11px,1vw,19px);font-weight:800;color:var(--blue);letter-spacing:.12em}}.cover-rule{{width:8%;height:7px;margin:4.5% 0 3%;background:var(--coral)}}.cover h1{{max-width:78%;margin:0;font-size:clamp(38px,5vw,82px);line-height:1.12;letter-spacing:0}}.cover p{{max-width:68%;margin:2.2% 0 0;color:var(--muted);font-size:clamp(16px,1.7vw,29px);line-height:1.5}}.cover-meta{{position:absolute;left:8.4%;bottom:8%;display:flex;gap:28px;color:var(--muted);font-size:clamp(10px,.85vw,15px)}}
.slide-header{{height:12%;display:grid;grid-template-columns:4.5% 1fr auto;align-items:start;gap:1.4%;border-bottom:2px solid var(--line)}}.slide-header span{{padding-top:4px;color:var(--coral);font-weight:800;font-size:clamp(12px,1vw,19px)}}.slide-header h2{{margin:0;font-size:clamp(24px,2.7vw,45px);line-height:1.1;letter-spacing:0;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}.slide-header em{{font-style:normal;font-weight:800;font-size:clamp(9px,.8vw,14px)}}.layout{{height:76%;padding-top:4.2%}}
.statement-layout{{display:grid;grid-template-columns:1.12fr .88fr;gap:7%;align-items:stretch}}.statement{{border-left:9px solid var(--coral);padding:4% 4% 3% 6%;display:flex;flex-direction:column;justify-content:center}}small,.insights h3{{color:var(--blue);font-weight:800;font-size:clamp(11px,1vw,18px)}}.statement h3,.editorial-lead h3{{margin:7% 0 0;font-size:clamp(26px,3.15vw,52px);line-height:1.27;letter-spacing:0}}.evidence,.insights ol,.editorial-layout ol{{list-style:none;margin:0;padding:0;display:grid;align-content:center}}.evidence li,.insights li,.editorial-layout li{{display:grid;grid-template-columns:11% 1fr;gap:4%;padding:6% 0;border-bottom:2px solid var(--line)}}.evidence span,.insights span,.editorial-layout span{{color:var(--coral);font-weight:800}}li p{{margin:0;font-size:clamp(15px,1.45vw,25px);line-height:1.45}}
.chart-layout{{display:grid;grid-template-columns:1.65fr .75fr;gap:5%}}.chart{{min-width:0;height:100%;background:#fff;border-left:7px solid var(--blue);padding:3%}}.chart svg{{width:100%;height:100%;overflow:visible}}.chart-title{{font-size:23px;font-weight:800;fill:var(--ink)}}.chart-label{{font-size:13px;fill:var(--muted)}}.insights{{padding-top:2%}}.insights h3{{margin:0 0 4%}}.insights li{{grid-template-columns:13% 1fr;padding:8% 0}}
.metric-layout{{display:grid;grid-template-columns:.9fr 1.1fr;gap:7%;align-items:center}}.metric{{height:88%;padding:7%;background:var(--ink);color:#fff;border-top:10px solid var(--coral);display:flex;flex-direction:column;justify-content:center}}.metric small{{color:#9fc5ff}}.metric strong{{margin:7% 0 5%;font-size:clamp(44px,6vw,96px);line-height:1;color:#fff}}.metric p{{margin:0;font-size:clamp(17px,1.7vw,28px);line-height:1.45}}.metric-layout ul{{margin:0;padding:0;list-style:none}}.metric-layout li{{padding:5% 0;border-bottom:2px solid var(--line);font-size:clamp(17px,1.6vw,27px);line-height:1.45}}
.editorial-layout{{display:grid;grid-template-columns:1fr 1fr;gap:8%}}.editorial-lead{{align-self:center;padding:4% 0}}.editorial-lead:after{{content:"";display:block;width:24%;height:7px;margin-top:9%;background:var(--teal)}}.editorial-layout ol{{align-content:center}}.action-layout{{position:relative;display:flex;align-items:center}}.action-line{{position:absolute;left:4%;right:4%;top:45%;height:4px;background:var(--line)}}.action-layout ol{{position:relative;z-index:1;width:100%;display:grid;grid-template-columns:repeat(4,1fr);gap:3%;margin:0;padding:0;list-style:none}}.action-layout li{{min-height:220px;padding:10% 8%;background:#fff;border-top:8px solid var(--blue)}}.action-layout span{{display:inline-grid;place-items:center;width:44px;height:44px;background:var(--coral);color:#fff;font-weight:800}}.action-layout p{{margin-top:16%}}
.references ol{{margin:4% 2%;padding-left:4%;columns:2;column-gap:7%}}.references li{{break-inside:avoid;margin:0 0 4%;padding-left:2%;color:#334155;font-size:clamp(12px,1.1vw,19px);line-height:1.45}}sup{{margin-left:.2em;color:var(--blue);font-size:.55em;vertical-align:super}}.slide-source{{position:absolute;left:6.6%;bottom:3.1%;max-width:76%;color:var(--muted);font-size:clamp(8px,.68vw,12px);white-space:nowrap;overflow:hidden;text-overflow:ellipsis}}.slide-number{{position:absolute;right:6.6%;bottom:3%;color:var(--muted);font-size:clamp(9px,.75vw,13px)}}
.controls{{position:absolute;z-index:20;left:50%;bottom:14px;transform:translateX(-50%);display:flex;align-items:center;gap:10px;padding:7px 10px;border:1px solid #ffffff2e;background:#111a27d9;color:#fff;backdrop-filter:blur(10px)}}.controls button{{width:32px;height:28px;border:0;background:transparent;color:#fff;font-size:18px;cursor:pointer}}.controls span{{min-width:64px;text-align:center;font-size:12px}}.progress{{position:absolute;z-index:21;left:0;bottom:0;height:4px;background:var(--coral);transition:width .3s ease}}.format-note{{position:absolute;z-index:20;right:14px;top:12px;padding:6px 9px;background:#111a27c9;color:#dce7f6;font-size:10px}}
@media(max-width:720px){{.controls{{bottom:7px}}.format-note{{display:none}}}}
@media(prefers-reduced-motion:reduce){{.slide.active .layout>*{{animation:none}}}}
</style></head><body><main class="deck"><div class="stage" id="stage">{slides}</div><div class="format-note">网页演示 · HTML + JS · 非 PowerPoint 文件</div><nav class="controls" aria-label="演示控制"><button id="prev" aria-label="上一页">‹</button><span id="counter"></span><button id="next" aria-label="下一页">›</button><button id="full" aria-label="全屏">⛶</button></nav><div class="progress" id="progress"></div></main>
<script type="application/json" id="chart-data">{chart_data}</script><script>
(() => {{
  const slides=[...document.querySelectorAll('.slide')], counter=document.querySelector('#counter'), progress=document.querySelector('#progress'); let index=0;
  const charts=JSON.parse(document.querySelector('#chart-data').textContent||'[]');
  const colors=['#1769e0','#0b9f91','#f05a47','#f2bd42'];
  function show(next){{index=(next+slides.length)%slides.length;slides.forEach((s,i)=>s.classList.toggle('active',i===index));counter.textContent=`${{index+1}} / ${{slides.length}}`;progress.style.width=`${{(index+1)/slides.length*100}}%`;}}
  function svg(tag,attrs={{}}){{const el=document.createElementNS('http://www.w3.org/2000/svg',tag);Object.entries(attrs).forEach(([k,v])=>el.setAttribute(k,String(v)));return el}}
  document.querySelectorAll('[data-chart]').forEach(host=>{{const c=charts[Number(host.dataset.chart)];if(!c)return;const title=document.createElement('div');title.style.cssText='font-size:20px;font-weight:800;margin:0 0 10px';title.textContent=c.title;host.append(title);if(c.type==='pie'){{const values=c.series[0]?.values||[],total=values.reduce((a,b)=>a+b,0)||1;let angle=0;const stops=values.map((v,i)=>{{const start=angle;angle+=v/total*360;return `${{colors[i%colors.length]}} ${{start}}deg ${{angle}}deg`}});const pie=document.createElement('div');pie.style.cssText=`width:min(54%,330px);aspect-ratio:1;border-radius:50%;margin:5% auto;background:conic-gradient(${{stops.join(',')}})`;host.append(pie);return}}
    const box=svg('svg',{{viewBox:'0 0 760 390',role:'img','aria-label':c.title}}), values=c.series.flatMap(s=>s.values), max=Math.max(...values,1), min=Math.min(0,...values), range=max-min||1;for(let i=0;i<5;i++){{const y=62+i*65;box.append(svg('line',{{x1:55,y1:y,x2:735,y2:y,stroke:'#dce2e9','stroke-width':1}}))}}const count=c.categories.length,slot=660/Math.max(1,count);c.categories.forEach((label,i)=>{{const t=svg('text',{{x:70+i*slot+slot/2,y:370,'text-anchor':'middle',class:'chart-label'}});t.textContent=label;box.append(t)}});c.series.forEach((series,si)=>{{if(c.type==='line'){{const pts=series.values.map((v,i)=>`${{70+i*slot+slot/2}},${{322-(v-min)/range*255}}`).join(' ');box.append(svg('polyline',{{points:pts,fill:'none',stroke:colors[si%colors.length],'stroke-width':5,'stroke-linejoin':'round'}}))}}else{{const bw=Math.min(42,slot*.66/c.series.length);series.values.forEach((v,i)=>{{const h=(v-min)/range*255,x=70+i*slot+slot/2+(si-(c.series.length-1)/2)*bw;box.append(svg('rect',{{x:x-bw*.42,y:322-h,width:bw*.84,height:h,fill:colors[si%colors.length]}}))}})}}}});host.append(box)}});
  document.querySelector('#prev').onclick=()=>show(index-1);document.querySelector('#next').onclick=()=>show(index+1);document.querySelector('#full').onclick=()=>document.documentElement.requestFullscreen?.();document.addEventListener('keydown',e=>{{if(['ArrowRight','PageDown',' '].includes(e.key))show(index+1);if(['ArrowLeft','PageUp'].includes(e.key))show(index-1);if(e.key==='Home')show(0);if(e.key==='End')show(slides.length-1)}});show(0);
}})();
</script><!-- Visual approach inspired by frontend-slides (MIT), Copyright (c) 2025 Zara Zhang. --></body></html>"""


def create_docx(request: DeliverableRequest) -> bytes:
    request = _normalize_document_request(request)
    document = Document()
    styles = document.styles
    for style_name in ("Normal", "Title", "Heading 1", "Heading 2"):
        style = styles[style_name]
        style.font.name = "Arial"
        style._element.get_or_add_rPr().get_or_add_rFonts().set(qn("w:eastAsia"), "PingFang SC")
    styles["Normal"].font.size = Pt(10.5)
    title = document.add_heading(request.title, 0)
    title.runs[0].font.color.rgb = RGBColor(31, 41, 55)
    if request.subtitle:
        paragraph = document.add_paragraph(request.subtitle)
        paragraph.runs[0].font.color.rgb = RGBColor(91, 105, 125)
    for section_index, section in enumerate(request.sections):
        if section_index and valid_chart(section.chart):
            document.add_page_break()
        heading = document.add_heading(section.heading, level=1)
        heading.paragraph_format.keep_with_next = True
        for content in section.paragraphs:
            document.add_paragraph(content)
        for bullet in section.bullets:
            document.add_paragraph(bullet, style="List Bullet")
        if valid_chart(section.chart):
            picture = io.BytesIO(chart_png(section.chart))
            document.add_picture(picture, width=DocxInches(6.35))
            source = normalize_markers(section.chart.source_ref, request) or inline_sources(request, section.refs)
            caption = document.add_paragraph("图表来源：" + source)
            caption.alignment = 1
            for run in caption.runs:
                run.font.size = Pt(8.5)
                run.font.color.rgb = RGBColor(91, 105, 125)
    entries = reference_entries(request)
    if entries:
        document.add_page_break()
        document.add_heading("参考文献", level=1)
        for entry in entries:
            paragraph = document.add_paragraph(entry)
            paragraph.paragraph_format.left_indent = DocxInches(0.25)
            paragraph.paragraph_format.first_line_indent = DocxInches(-0.25)
            for run in paragraph.runs:
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor(71, 85, 105)
    _apply_docx_fonts(document)
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def create_pdf(request: DeliverableRequest) -> bytes:
    request = _normalize_document_request(request)
    stream = io.BytesIO()
    font_name = _pdf_font()
    document = SimpleDocTemplate(
        stream,
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
        title=request.title,
        author="FinFlow Studio",
    )
    base = getSampleStyleSheet()
    title_style = ParagraphStyle("FinFlowTitle", parent=base["Title"], fontName=font_name,
                                 fontSize=24, leading=32, textColor=colors.HexColor("#1f2937"),
                                 alignment=TA_LEFT, spaceAfter=8 * mm)
    subtitle_style = ParagraphStyle("FinFlowSubtitle", parent=base["Normal"], fontName=font_name,
                                    fontSize=11, leading=18, textColor=colors.HexColor("#5b697d"),
                                    spaceAfter=9 * mm)
    heading_style = ParagraphStyle("FinFlowHeading", parent=base["Heading1"], fontName=font_name,
                                   fontSize=16, leading=24, textColor=colors.HexColor("#19487e"),
                                   spaceBefore=5 * mm, spaceAfter=4 * mm)
    body_style = ParagraphStyle("FinFlowBody", parent=base["BodyText"], fontName=font_name,
                                fontSize=10.5, leading=18, textColor=colors.HexColor("#1f2937"),
                                spaceAfter=3 * mm)
    ref_style = ParagraphStyle("FinFlowRef", parent=body_style, fontSize=8.5, leading=13,
                               textColor=colors.HexColor("#5b697d"))
    story = [Paragraph(_pdf_escape(request.title), title_style)]
    if request.subtitle:
        story.extend([Paragraph(_pdf_escape(request.subtitle), subtitle_style), Spacer(1, 2 * mm)])
    for section_index, section in enumerate(request.sections):
        if section_index and valid_chart(section.chart):
            story.append(PageBreak())
        story.append(Paragraph(_pdf_escape(section.heading), heading_style))
        story.extend(Paragraph(_pdf_escape(content), body_style) for content in section.paragraphs if content.strip())
        if section.bullets:
            story.append(ListFlowable(
                [ListItem(Paragraph(_pdf_escape(item), body_style), leftIndent=5 * mm) for item in section.bullets if item.strip()],
                bulletType="bullet", leftIndent=7 * mm, bulletFontName=font_name,
            ))
            story.append(Spacer(1, 3 * mm))
        if valid_chart(section.chart):
            chart_stream = io.BytesIO(chart_png(section.chart))
            image = ReportLabImage(chart_stream, width=165 * mm, height=89 * mm)
            image.hAlign = "CENTER"
            story.extend([image, Spacer(1, 2 * mm)])
            source = normalize_markers(section.chart.source_ref, request) or inline_sources(request, section.refs)
            if source:
                story.append(Paragraph(_pdf_escape("图表来源：" + source), ref_style))
    entries = reference_entries(request)
    if entries:
        story.extend([PageBreak(), Paragraph(_pdf_escape("参考文献"), heading_style)])
        story.extend(Paragraph(_pdf_escape(entry), ref_style) for entry in entries)
    document.build(story, onFirstPage=_pdf_page, onLaterPages=_pdf_page)
    return stream.getvalue()


def create_financial_report(request: DeliverableRequest) -> bytes:
    request = _normalize_document_request(request)
    specification = {
        "schema_version": 2,
        "renderer": "finflow-echarts-perspective",
        "title": request.title,
        "subtitle": request.subtitle or "由 FinFlow Studio 工作流生成",
        "theme": request.theme,
        "sections": [section.model_dump(mode="json") for section in request.sections],
        "references": reference_records(request),
    }
    return json.dumps(specification, ensure_ascii=False, separators=(",", ":")).encode("utf-8")


def create_mermaid(request: DeliverableRequest) -> bytes:
    generated = _mermaid_source(request)
    if generated:
        return (generated.rstrip() + "\n").encode("utf-8")
    lines = ["flowchart TD", '  root["%s"]' % _escape(request.title)]
    for index, section in enumerate(request.sections, start=1):
        section_id = "section_%d" % index
        lines.append('  %s["%s"]' % (section_id, _escape(section.heading)))
        lines.append("  root --> %s" % section_id)
        contents = list(section.bullets) or list(section.paragraphs)
        for item_index, item in enumerate(contents[:20], start=1):
            item_id = "%s_item_%d" % (section_id, item_index)
            lines.append('  %s["%s"]' % (item_id, _escape(item)))
            lines.append("  %s --> %s" % (section_id, item_id))
    return ("\n".join(lines) + "\n").encode("utf-8")


def create_excalidraw(request: DeliverableRequest) -> bytes:
    source = _mermaid_source(request) or create_mermaid(request).decode("utf-8")
    direction_match = re.search(r"^(?:flowchart|graph)\s+(TD|TB|BT|RL|LR)", source, re.IGNORECASE)
    direction = direction_match.group(1).upper() if direction_match else "TD"
    node_pattern = re.compile(
        r"\b([A-Za-z][A-Za-z0-9_]*)\s*(?:\[\"([^\"]+)\"\]|\[([^\]]+)\]|\{([^}]+)\}|\(([^)]+)\))"
    )
    labels: dict[str, str] = {}
    for match in node_pattern.finditer(source):
        labels.setdefault(match.group(1), next(value for value in match.groups()[1:] if value is not None).strip())
    edge_pattern = re.compile(r"\b([A-Za-z][A-Za-z0-9_]*)\s*-->(?:\|([^|]+)\|)?\s*([A-Za-z][A-Za-z0-9_]*)")
    edges = [(match.group(1), match.group(3), (match.group(2) or "").strip()) for match in edge_pattern.finditer(source)]
    for source_id, target_id, _ in edges:
        labels.setdefault(source_id, source_id)
        labels.setdefault(target_id, target_id)
    if not labels:
        points = [point for section in request.sections for point in _content_points(section)][:12]
        labels = {f"node{index}": point for index, point in enumerate(points, start=1)}
        keys = list(labels)
        edges = [(keys[index], keys[index + 1], "") for index in range(len(keys) - 1)]

    positions: dict[str, tuple[float, float]] = {}
    horizontal = direction in {"LR", "RL"}
    for index, node_id in enumerate(labels):
        primary = index * 260
        x, y = (80 + primary, 120 + (index % 2) * 150) if horizontal else (100 + (index % 3) * 280, 80 + (index // 3) * 180)
        positions[node_id] = (x, y)

    elements: list[dict] = []
    for index, (node_id, label) in enumerate(labels.items(), start=1):
        x, y = positions[node_id]
        rect_id = "shape_" + node_id
        text_id = "text_" + node_id
        elements.append(_excalidraw_element(rect_id, "rectangle", x, y, 210, 84, index, {
            "backgroundColor": "#e7f1ff", "fillStyle": "solid", "roundness": {"type": 3},
            "boundElements": [{"type": "text", "id": text_id}],
        }))
        text = label[:100]
        elements.append(_excalidraw_element(text_id, "text", x + 14, y + 24, 182, 36, index + 1000, {
            "strokeColor": "#1f2937", "fontSize": 18, "fontFamily": 5, "text": text,
            "textAlign": "center", "verticalAlign": "middle", "containerId": rect_id,
            "originalText": text, "autoResize": True, "lineHeight": 1.25,
        }))
    for index, (source_id, target_id, label) in enumerate(edges, start=1):
        if source_id not in positions or target_id not in positions:
            continue
        sx, sy = positions[source_id][0] + 105, positions[source_id][1] + 42
        tx, ty = positions[target_id][0] + 105, positions[target_id][1] + 42
        arrow_id = f"arrow_{index}"
        elements.append(_excalidraw_element(arrow_id, "arrow", sx, sy, abs(tx - sx), abs(ty - sy), index + 2000, {
            "points": [[0, 0], [tx - sx, ty - sy]], "lastCommittedPoint": None,
            "startBinding": {"elementId": "shape_" + source_id, "focus": 0, "gap": 6},
            "endBinding": {"elementId": "shape_" + target_id, "focus": 0, "gap": 6},
            "startArrowhead": None, "endArrowhead": "arrow", "elbowed": False,
        }))
        if label:
            mid_x, mid_y = (sx + tx) / 2, (sy + ty) / 2
            elements.append(_excalidraw_element(f"edge_text_{index}", "text", mid_x - 40, mid_y - 22, 80, 22, index + 3000, {
                "strokeColor": "#5b697d", "fontSize": 14, "fontFamily": 5, "text": label[:40],
                "textAlign": "center", "verticalAlign": "middle", "containerId": None,
                "originalText": label[:40], "autoResize": True, "lineHeight": 1.25,
            }))
    payload = {
        "type": "excalidraw", "version": 2, "source": "https://finflow.local",
        "elements": elements,
        "appState": {"gridSize": None, "viewBackgroundColor": "#ffffff"}, "files": {},
    }
    return json.dumps(payload, ensure_ascii=False, separators=(",", ":")).encode("utf-8")


def _excalidraw_element(element_id: str, element_type: str, x: float, y: float, width: float,
                        height: float, seed: int, extra: dict) -> dict:
    item = {
        "id": element_id, "type": element_type, "x": x, "y": y, "width": width, "height": height,
        "angle": 0, "strokeColor": "#236dcc", "backgroundColor": "transparent", "fillStyle": "solid",
        "strokeWidth": 2, "strokeStyle": "solid", "roughness": 2, "opacity": 100, "groupIds": [],
        "frameId": None, "roundness": None, "seed": seed * 7919, "version": 1,
        "versionNonce": seed * 104729, "isDeleted": False, "boundElements": [], "updated": 1,
        "link": None, "locked": False,
    }
    item.update(extra)
    return item


def _mermaid_source(request: DeliverableRequest) -> str:
    for section in request.sections:
        for value in section.paragraphs:
            clean = re.sub(r"^```(?:mermaid)?\s*|\s*```$", "", value.strip(), flags=re.IGNORECASE)
            if re.match(r"^(flowchart|graph)\s+(TD|TB|BT|RL|LR)\b", clean, flags=re.IGNORECASE):
                return clean
    return ""


def _add_text(slide, text: str, left: float, top: float, width: float, height: float,
              size: int, color: PptColor, bold: bool, min_size: int = 12,
              single_line: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    normalized = re.sub(r"\s*\n+\s*", " " if single_line else "\n", str(text)).strip()
    fitted_size = _ppt_fit_font_size(normalized, width, height, size, min_size, single_line)
    for index, line in enumerate(normalized.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.08
        paragraph.space_after = PptPt(0)
        parts = re.split(r"(\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\))", line)
        for part in parts:
            if not part:
                continue
            run = paragraph.add_run()
            run.text = part
            citation = bool(re.fullmatch(r"\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\)", part))
            run.font.size = PptPt(max(7, round(fitted_size * 0.58))) if citation else PptPt(fitted_size)
            run.font.bold = bold
            run.font.color.rgb = color
            _set_ppt_run_font(run, PPT_FONT)


def _ppt_fit_font_size(text: str, width: float, height: float, preferred: int, minimum: int,
                       single_line: bool) -> int:
    units = sum(1.0 if ord(char) > 255 else 0.56 for char in re.sub(r"\s+", " ", text))
    explicit_lines = max(1, text.count("\n") + 1)
    for candidate in range(preferred, minimum - 1, -1):
        capacity = max(1.0, width * 72 / (candidate * 0.94))
        lines = explicit_lines if single_line else max(explicit_lines, math.ceil(units / capacity))
        if single_line and units > capacity:
            continue
        if lines * candidate * 1.18 <= height * 72:
            return candidate
    return minimum


def _set_ppt_run_font(run, family: str) -> None:
    run.font.name = family
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        element = properties.find(pptx_qn(tag))
        if element is None:
            element = SubElement(properties, pptx_qn(tag))
        element.set("typeface", family)


def _add_reference_slides(presentation: Presentation, request: DeliverableRequest, page_number: int) -> None:
    entries = reference_entries(request)
    for page_index, group in enumerate(_chunks(entries, 8), start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        heading = "参考文献" if page_index == 1 else "参考文献（续）"
        _add_slide_heading(slide, heading, page_number)
        for index, entry in enumerate(group):
            _add_text(slide, entry, 1.0, 1.65 + index * 0.62, 11.2, 0.48, 11, DARK, False)
        _add_brand(slide, page_number)
        page_number += 1


def _add_brand(slide, page_number: int) -> None:
    _add_text(slide, "FINFLOW STUDIO", 0.85, 0.55, 3.0, 0.25, 10, BLUE, True)
    _add_text(slide, str(page_number), 12.0, 6.9, 0.45, 0.2, 9, MUTED, False)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.55), Inches(1.05), Inches(0.07))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()


def _add_slide_heading(slide, heading: str, page_number: int) -> None:
    _add_text(slide, heading, 0.85, 0.72, 11.0, 0.62, 32, DARK, True,
              min_size=24, single_line=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.43), Inches(11.65), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_LINE
    line.line.fill.background()


def _apply_docx_fonts(document: Document) -> None:
    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            run.font.name = "PingFang SC"
            fonts = run._element.get_or_add_rPr().get_or_add_rFonts()
            for key in ("ascii", "hAnsi", "eastAsia", "cs"):
                fonts.set(qn("w:" + key), "PingFang SC")


def _add_point_list(slide, points: List[str]) -> None:
    top = 1.72
    height = 1.12 if len(points) >= 4 else 1.35
    for index, point in enumerate(points, start=1):
        _add_text(slide, point, 1.65, top, 10.65, height, 18, DARK, False)
        number = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(top + 0.05), Inches(0.5), Inches(0.5))
        number.fill.solid()
        number.fill.fore_color.rgb = PALE_BLUE
        number.line.color.rgb = LIGHT_LINE
        frame = number.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = f"{index:02d}"
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = PptPt(13)
        paragraph.font.bold = True
        paragraph.font.color.rgb = DARK_BLUE
        top += height + 0.12


def _add_compact_points(slide, points: List[str], left: float, top: float, width: float) -> None:
    for index, point in enumerate(points, start=1):
        _add_text(slide, f"{index:02d}", left, top + 0.03, 0.42, 0.3, 11, DARK_BLUE, True)
        _add_text(slide, point, left + 0.52, top, width - 0.52, 1.12, 15, DARK, False)
        top += 1.35


def _content_points(section: DeliverableSection) -> List[str]:
    raw: List[str] = []
    for paragraph in section.paragraphs:
        raw.extend(re.split(r"\n+|(?<=[。！？；])", paragraph))
    raw.extend(section.bullets)
    points: List[str] = []
    seen = set()
    for item in raw:
        clean = _clean_ppt_text(item)
        if len(clean) < 4:
            continue
        for part in _wrap_point(clean, 72):
            key = part.lower()
            if key not in seen:
                points.append(part)
                seen.add(key)
            if len(points) >= 36:
                return points
    return points or ["工作流已完成处理，但当前步骤没有生成可展示的文字内容。"]


def _normalize_ppt_request(request: DeliverableRequest) -> DeliverableRequest:
    normalized: List[DeliverableSection] = []
    for section in request.sections:
        generated = normalize_markers("\n".join(section.paragraphs).strip(), request)
        slides = _json_ppt_sections(generated, section) or _labeled_ppt_sections(generated, section)
        normalized.extend(slides or [section])
    return request.model_copy(update={"sections": normalized}, deep=True)


def _normalize_document_request(request: DeliverableRequest) -> DeliverableRequest:
    normalized: List[DeliverableSection] = []
    for section in request.sections:
        generated = normalize_markers("\n".join(section.paragraphs).strip(), request)
        normalized.extend(_json_report_sections(generated, section) or [section])
    return request.model_copy(update={"sections": normalized}, deep=True)


def _json_ppt_sections(value: str, source: DeliverableSection) -> List[DeliverableSection]:
    candidate = re.sub(r"^```(?:json)?\s*", "", value.strip(), flags=re.IGNORECASE)
    candidate = re.sub(r"\s*```$", "", candidate)
    try:
        payload = json.loads(candidate)
    except (json.JSONDecodeError, TypeError):
        return []
    slides = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(slides, list):
        return []
    result: List[DeliverableSection] = []
    for slide in slides[:12]:
        if not isinstance(slide, dict):
            continue
        heading = _clean_slide_heading(str(slide.get("title") or ""))
        summary = _compact_ppt_copy(str(slide.get("summary") or ""), 48)
        bullets = slide.get("bullets")
        clean_bullets = [_compact_ppt_copy(str(item), 42) for item in bullets] if isinstance(bullets, list) else []
        bullet_limit = 3 if summary else 4
        clean_bullets = [item for item in clean_bullets if len(item) >= 4][:bullet_limit]
        if heading and (summary or clean_bullets):
            result.append(DeliverableSection(
                heading=heading,
                paragraphs=[summary] if summary else [],
                bullets=clean_bullets,
                refs=source.refs,
                chart=_parse_chart(slide.get("chart")),
            ))
    return result


def _json_report_sections(value: str, source: DeliverableSection) -> List[DeliverableSection]:
    candidate = re.sub(r"^```(?:json)?\s*", "", value.strip(), flags=re.IGNORECASE)
    candidate = re.sub(r"\s*```$", "", candidate)
    try:
        payload = json.loads(candidate)
    except (json.JSONDecodeError, TypeError):
        return []
    items = payload.get("sections") if isinstance(payload, dict) else None
    if not isinstance(items, list):
        items = payload.get("slides") if isinstance(payload, dict) else None
    if not isinstance(items, list):
        return []
    result: List[DeliverableSection] = []
    for item in items[:30]:
        if not isinstance(item, dict):
            continue
        heading = _clean_slide_heading(str(item.get("heading") or item.get("title") or ""))
        paragraphs = item.get("paragraphs")
        if not isinstance(paragraphs, list):
            summary = _clean_ppt_text(str(item.get("summary") or ""))
            paragraphs = [summary] if summary else []
        bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
        clean_paragraphs = [_clean_ppt_text(str(value)) for value in paragraphs if str(value).strip()][:8]
        clean_bullets = [_clean_ppt_text(str(value)) for value in bullets if str(value).strip()][:12]
        if heading and (clean_paragraphs or clean_bullets or _parse_chart(item.get("chart"))):
            result.append(DeliverableSection(
                heading=heading, paragraphs=clean_paragraphs, bullets=clean_bullets,
                refs=source.refs, chart=_parse_chart(item.get("chart")),
            ))
    return result


def _parse_chart(value: object) -> Optional[DeliverableChart]:
    if not isinstance(value, dict):
        return None
    try:
        chart = DeliverableChart.model_validate(value)
        return chart if valid_chart(chart) else None
    except (TypeError, ValueError):
        return None


def _labeled_ppt_sections(value: str, source: DeliverableSection) -> List[DeliverableSection]:
    lines = [item.strip() for item in re.split(r"\n+", value) if item.strip()]
    result: List[DeliverableSection] = []
    heading = ""
    points: List[str] = []
    page_pattern = re.compile(r"^第\s*[一二三四五六七八九十百零〇\d]+\s*页\s*[｜|:：\-—]?\s*(.*)$")

    def append_page() -> None:
        clean_points = [_clean_ppt_text(item) for item in points]
        clean_points[:] = [item for item in clean_points if len(item) >= 4][:4]
        if heading and clean_points:
            result.append(DeliverableSection(
                heading=_clean_slide_heading(heading),
                paragraphs=clean_points[:1],
                bullets=clean_points[1:5],
                refs=source.refs,
            ))

    for line in lines:
        match = page_pattern.match(_clean_ppt_text(line))
        if match:
            append_page()
            heading = match.group(1).strip() or source.heading
            points = []
        elif heading:
            points.extend(part for part in re.split(r"(?<=[。！？；])", line) if part.strip())
    append_page()
    return result if len(result) >= 2 else []


def _clean_slide_heading(value: str) -> str:
    clean = _clean_ppt_text(value)
    clean = re.sub(r"^(?:第\s*[一二三四五六七八九十百零〇\d]+\s*页|slide\s*\d+)\s*[｜|:：\-—]?\s*", "", clean, flags=re.IGNORECASE)
    replacements = {
        "管理层需要关注的": "",
        "仍需同步关注": "需关注",
        "呈现先承压后快速修复的经营轨迹": "先承压后快速修复",
        "的主要判断": "",
        "的核心结论": "",
    }
    for source, target in replacements.items():
        clean = clean.replace(source, target)
    if len(clean) <= 26:
        return clean
    clauses = [item.strip() for item in re.split(r"[，,；;。:：]", clean) if item.strip()]
    concise = clauses[0] if clauses else clean
    return concise[:26]


def _compact_ppt_copy(value: str, limit: int) -> str:
    clean = _clean_ppt_text(value)
    if len(clean) <= limit:
        return clean
    markers = "".join(re.findall(r"\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\)", clean))
    body = re.sub(r"\s*(?:\[\d+(?:[-,]\d+)*\]|\([^)]*n\.d\.[^)]*\))", "", clean).strip()
    clauses = [item.strip(" ，,；;。") for item in re.split(r"[；;。]|(?<=，)|(?<=,)", body) if item.strip()]
    compact = ""
    for clause in clauses:
        candidate = (compact + clause).strip()
        if compact and len(candidate) > limit - len(markers):
            break
        compact = candidate
    if not compact:
        compact = body[:max(1, limit - len(markers))].rstrip("，,；;。")
    return (compact + (" " + markers if markers else ""))[:limit]


def _clean_ppt_text(value: str) -> str:
    clean = value.replace("**", "").replace("`", "")
    clean = re.sub(r"^\s*(?:[#>*•\-]+\s*|\d+[.、]\s+)+", "", clean)
    return re.sub(r"\s+", " ", clean).strip(" ；;。")


def _wrap_point(value: str, limit: int) -> List[str]:
    if len(value) <= limit:
        return [value]
    parts = [part.strip() for part in re.split(r"(?<=[，,；;。])", value) if part.strip()]
    result: List[str] = []
    current = ""
    for part in parts:
        if current and len(current) + len(part) > limit:
            result.append(current[:limit])
            current = part
        else:
            current += part
    if current:
        result.append(current[:limit])
    return result


def _chunks(values: List[str], size: int) -> Iterable[List[str]]:
    for index in range(0, len(values), size):
        yield values[index:index + size]


def _format_refs(section: DeliverableSection) -> str:
    values = []
    for ref in section.refs:
        location = "、".join("%s=%s" % (key, value) for key, value in ref.location.items())
        values.append(ref.source_name + ("（" + location + "）" if location else ""))
    return "；".join(values)


def _pdf_font() -> str:
    name = "FinFlowCJK"
    if name in pdfmetrics.getRegisteredFontNames():
        return name
    candidates = [
        Path("/System/Library/Fonts/Supplemental/Arial Unicode.ttf"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
    ]
    for candidate in candidates:
        if candidate.exists():
            pdfmetrics.registerFont(TTFont(name, str(candidate)))
            return name
    return "Helvetica"


def _pdf_page(canvas, document) -> None:
    canvas.saveState()
    canvas.setStrokeColor(colors.HexColor("#dbe5f1"))
    canvas.line(20 * mm, 14 * mm, A4[0] - 20 * mm, 14 * mm)
    canvas.setFillColor(colors.HexColor("#5b697d"))
    canvas.setFont(_pdf_font(), 8)
    canvas.drawString(20 * mm, 9 * mm, "FINFLOW STUDIO")
    canvas.drawRightString(A4[0] - 20 * mm, 9 * mm, str(document.page))
    canvas.restoreState()


def _pdf_escape(value: str) -> str:
    return value.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")


def _escape(value: str) -> str:
    return re.sub(r"[\"\[\]{}]", " ", value).replace("\n", " ")[:160]
