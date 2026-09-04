import json
from pathlib import Path
from typing import List

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from .models import DocumentPreview, PreviewBlock, PreviewPage
from .research import _json_tables


MAX_PAGES = 300
MAX_BLOCKS = 3000
MAX_TABLE_ROWS = 200
MAX_TEXT_CHARS = 20_000
MAX_DIAGRAM_CHARS = 5_000_000


def preview_document(path: Path, file_name: str) -> DocumentPreview:
    suffix = Path(file_name).suffix.lower()
    if suffix == ".pptx":
        return _preview_pptx(path, file_name)
    if suffix == ".docx":
        return _preview_docx(path, file_name)
    if suffix == ".pdf":
        return _preview_pdf(path, file_name)
    if suffix == ".json":
        payload = json.loads(path.read_text(encoding="utf-8-sig", errors="replace"))
        declared_tables = payload.get("tables") if isinstance(payload, dict) else payload
        blocks: list[PreviewBlock] = []
        if isinstance(declared_tables, list) and declared_tables and all(
                isinstance(table, dict) and isinstance(table.get("rows"), list) for table in declared_tables):
            for table in declared_tables:
                blocks.append(PreviewBlock(type="heading", text=str(table.get("title", "数据表"))))
                blocks.append(PreviewBlock(type="table", rows=[
                    [str(cell) for cell in row] for row in table["rows"] if isinstance(row, list)
                ]))
        else:
            blocks = [PreviewBlock(type="table", rows=[[str(cell) for cell in row] for row in table["rows"]])
                      for table in _json_tables(payload)]
        if not blocks:
            blocks = [PreviewBlock(type="text", text=json.dumps(payload, ensure_ascii=False, indent=2)[:MAX_TEXT_CHARS])]
        return DocumentPreview(file_name=file_name, kind="data", title=Path(file_name).stem,
                               pages=[PreviewPage(number=1, title="", blocks=blocks)])
    if suffix in {".txt", ".md", ".mermaid", ".mmd", ".excalidraw"}:
        limit = MAX_DIAGRAM_CHARS if suffix in {".mermaid", ".mmd", ".excalidraw"} else MAX_TEXT_CHARS
        text = path.read_text(encoding="utf-8-sig", errors="replace")[:limit]
        return DocumentPreview(
            file_name=file_name,
            kind="text",
            title=Path(file_name).stem,
            pages=[PreviewPage(number=1, title="", blocks=[PreviewBlock(type="text", text=text)])],
        )
    raise ValueError("该文件请使用 PDF 或表格预览方式")


def _preview_pdf(path: Path, file_name: str) -> DocumentPreview:
    reader = PdfReader(str(path))
    pages: List[PreviewPage] = []
    warnings: List[str] = []
    for number, page in enumerate(reader.pages[:MAX_PAGES], start=1):
        text = (page.extract_text() or "").strip()[:MAX_TEXT_CHARS]
        pages.append(PreviewPage(number=number, title="", blocks=[PreviewBlock(type="text", text=text)]))
    if len(reader.pages) > MAX_PAGES:
        warnings.append("PDF 页数较多，在线预览仅显示前 %d 页" % MAX_PAGES)
    return DocumentPreview(file_name=file_name, kind="document", title=Path(file_name).stem,
                           pages=pages, warnings=warnings)


def _preview_pptx(path: Path, file_name: str) -> DocumentPreview:
    presentation = Presentation(str(path))
    pages: List[PreviewPage] = []
    warnings: List[str] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > MAX_PAGES:
            warnings.append("幻灯片较多，在线预览仅显示前 %d 页" % MAX_PAGES)
            break
        texts: List[str] = []
        for shape in slide.shapes:
            value = getattr(shape, "text", "").strip()
            if value:
                texts.append(value[:MAX_TEXT_CHARS])
        title = texts[0] if texts else "幻灯片 %d" % slide_number
        blocks = [PreviewBlock(type="text", text=value) for value in texts[1:]]
        pages.append(PreviewPage(number=slide_number, title=title, blocks=blocks))
    return DocumentPreview(
        file_name=file_name,
        kind="presentation",
        title=Path(file_name).stem,
        pages=pages,
        warnings=warnings,
    )


def _preview_docx(path: Path, file_name: str) -> DocumentPreview:
    document = Document(str(path))
    blocks: List[PreviewBlock] = []
    warnings: List[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (paragraph.style.name if paragraph.style else "").lower()
        block_type = "heading" if "heading" in style_name or "title" in style_name else "text"
        blocks.append(PreviewBlock(type=block_type, text=text[:MAX_TEXT_CHARS]))
        if len(blocks) >= MAX_BLOCKS:
            warnings.append("文档内容较多，在线预览已截断")
            break
    if len(blocks) < MAX_BLOCKS:
        for table in document.tables:
            rows = [[cell.text.strip()[:2000] for cell in row.cells] for row in table.rows[:MAX_TABLE_ROWS]]
            blocks.append(PreviewBlock(type="table", rows=rows))
            if len(table.rows) > MAX_TABLE_ROWS:
                warnings.append("文档表格较长，在线预览仅显示前 %d 行" % MAX_TABLE_ROWS)
    return DocumentPreview(
        file_name=file_name,
        kind="document",
        title=Path(file_name).stem,
        pages=[PreviewPage(number=1, title="", blocks=blocks)],
        warnings=warnings,
    )
