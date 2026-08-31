from pathlib import Path

from docx import Document
from pptx import Presentation

from app.document_parser import parse_document
from app.models import FormulaColumn, SpreadsheetTransformRequest
from app.spreadsheet_files import transform_spreadsheet


def test_docx_preserves_paragraph_location(tmp_path: Path) -> None:
    path = tmp_path / "经营报告.docx"
    document = Document()
    document.add_paragraph("本月收入同比增长 12%。")
    document.add_paragraph("毛利率下降主要受到原材料成本影响。")
    document.save(path)

    parsed = parse_document(path, path.name, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    assert parsed.text_length > 10
    assert parsed.chunks[0].location["paragraph"] == 1
    assert parsed.chunks[0].content_hash.startswith("sha256:")


def test_pptx_preserves_slide_location(tmp_path: Path) -> None:
    path = tmp_path / "复盘.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "经营复盘"
    slide.placeholders[1].text = "收入达成预算，费用率需要继续关注。"
    presentation.save(path)

    parsed = parse_document(path, path.name, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    assert parsed.chunks[0].location["slide"] == 1
    assert "经营复盘" in parsed.chunks[0].text


def test_transform_workbook_preserves_formulas_and_creates_new_column(tmp_path: Path) -> None:
    import io
    from openpyxl import Workbook, load_workbook

    path = tmp_path / "finance.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["收入", "成本"])
    sheet.append([100, 40])
    sheet.append([100, 40])
    workbook.save(path)
    content = transform_spreadsheet(path, path.name, SpreadsheetTransformRequest(
        rename_headers={"收入": "营业收入"},
        formula_columns=[FormulaColumn(name="利润", formula="=A{row}-B{row}")],
        remove_duplicates=True,
    ))
    result = load_workbook(io.BytesIO(content), data_only=False)
    assert result.active.max_row == 2
    assert result.active["A1"].value == "营业收入"
    assert result.active["C2"].value == "=A2-B2"
    result.close()
