from io import BytesIO
import json

from docx import Document
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.util import Inches
from pypdf import PdfReader

from app.deliverables import create_docx, create_excalidraw, create_financial_report, create_html_slides, create_mermaid, create_pdf, create_pptx
from app.models import DeliverableRequest, DeliverableSection


def request() -> DeliverableRequest:
    return DeliverableRequest.model_validate({
        "title": "8 月经营复盘",
        "subtitle": "自动生成草稿",
        "sections": [{
            "heading": "核心结论",
            "bullets": ["收入同比增长 12%", "毛利率需要关注"],
            "refs": [{"source_name": "经营报告.pdf", "location": {"page": 12}}],
        }],
    })


def test_generates_real_office_files_with_refs() -> None:
    pptx = create_pptx(request())
    docx = create_docx(request())

    deck = Presentation(BytesIO(pptx))
    assert len(deck.slides) == 3
    assert "参考文献" in "\n".join(shape.text for shape in deck.slides[-1].shapes if hasattr(shape, "text"))
    document = Document(BytesIO(docx))
    assert "经营报告.pdf" in "\n".join(paragraph.text for paragraph in document.paragraphs)


def test_exact_total_slide_count_includes_cover() -> None:
    slides = [{"title": f"第{index}部分", "summary": "核心判断", "bullets": ["支撑信息"], "chart": None}
              for index in range(1, 13)]
    item = request().model_copy(update={
        "include_citations": False,
        "sections": [DeliverableSection(
            heading="分析结果",
            paragraphs=[json.dumps({"slides": slides, "total_slides": 12}, ensure_ascii=False)],
        )],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))

    assert len(deck.slides) == 12


def test_generates_readable_pdf() -> None:
    pdf = create_pdf(request())
    reader = PdfReader(BytesIO(pdf))

    assert pdf.startswith(b"%PDF")
    assert len(reader.pages) >= 1
    assert "8 月经营复盘" in (reader.pages[0].extract_text() or "")


def test_can_disable_all_citation_output() -> None:
    item = request().model_copy(update={"include_citations": False}, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert len(deck.slides) == 2
    assert "经营报告.pdf" not in text
    assert "参考文献" not in text


def test_apa_7_formats_reference_entries_without_ref_labels() -> None:
    plan = """{"slides":[{"title":"经营表现改善","summary":"收入保持增长 [Ref 1]","bullets":[],"chart":null}]}"""
    item = request().model_copy(update={
        "citation_style": "APA_7",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], refs=request().sections[0].refs)],
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    slide_text = "\n".join(shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text"))

    assert "(经营报告, n.d.)" in slide_text
    assert "(经营报告.pdf, n.d.)" not in slide_text
    assert "经营报告.pdf. (n.d.)." in text
    assert "[Ref" not in text


def test_apa_citation_compaction_tolerates_missing_comma() -> None:
    refs = [
        request().sections[0].refs[0].model_copy(
            update={"ref_id": "cost", "source_name": "device_cost_scenarios.csv"}
        ),
        request().sections[0].refs[0].model_copy(
            update={"ref_id": "rule", "source_name": "device_selection_rules.md"}
        ),
    ]
    plan = """{"slides":[{"title":"采购判断","summary":"成本占比超过八成（device_cost_scenarios.csvn.d.）","bullets":["完成验证（device_selection_rules.md, n.d.）"],"chart":null}]}"""
    item = request().model_copy(update={
        "citation_style": "APA_7",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], refs=refs)],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in list(deck.slides)[1:-1] for shape in slide.shapes if hasattr(shape, "text"))

    assert "device_cost_scenarios.csvn.d." not in text
    assert "device_selection_rules.md" not in text
    assert "(device cost, n.d.)" in text
    assert "(device selection, n.d.)" in text


def test_structured_ppt_reference_pages_only_include_used_sources() -> None:
    refs = [
        request().sections[0].refs[0],
        request().sections[0].refs[0].model_copy(update={"ref_id": "second", "source_name": "未使用的资料.pdf"}),
    ]
    plan = """{"slides":[{"title":"经营表现改善","summary":"收入保持增长 [Ref 1]","bullets":["数据已经核对 [Ref 1]"],"chart":null}]}"""
    item = request().model_copy(update={
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], refs=refs)],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert "经营报告.pdf" in text
    assert "未使用的资料.pdf" not in text


def test_generates_financial_report_workspace() -> None:
    report = json.loads(create_financial_report(request()))

    assert report["schema_version"] == 2
    assert report["renderer"] == "finbtp-echarts-perspective"
    assert report["title"] == "8 月经营复盘"
    assert report["sections"][0]["heading"] == "核心结论"
    assert report["references"][0]["source_name"] == "经营报告.pdf"
    assert report["references"][0]["location"] == {"page": 12}
    assert "经营报告.pdf" in report["references"][0]["formatted"]


def test_splits_long_analysis_into_readable_slides() -> None:
    item = request().model_copy(deep=True)
    item.sections[0].paragraphs = ["。".join(f"这是第 {index} 条需要在汇报中说明的业务结论" for index in range(12))]
    deck = Presentation(BytesIO(create_pptx(item)))

    assert len(deck.slides) >= 3
    assert all(len(shape.text) < 500 for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text_frame"))


def test_supports_management_decks_up_to_ten_slides() -> None:
    item = request().model_copy(update={"include_citations": False}, deep=True)
    item.sections[0].paragraphs = ["。".join(f"这是第 {index} 条管理层结论" for index in range(40))]
    deck = Presentation(BytesIO(create_pptx(item)))

    assert len(deck.slides) == 10


def test_huawei_style_c_skill_changes_native_pptx_design() -> None:
    item = request().model_copy(update={"ppt_skill": "guizang-huawei-style-c"}, deep=True)
    content = create_pptx(item)
    deck = Presentation(BytesIO(content))
    cover_text = "\n".join(shape.text for shape in deck.slides[0].shapes if hasattr(shape, "text_frame"))

    assert "企业汇报 · Classic Red" in cover_text
    assert any(
        shape.fill.type == MSO_FILL_TYPE.SOLID and str(shape.fill.fore_color.rgb) == "C7000B"
        for shape in deck.slides[0].shapes
    )


def test_generates_self_contained_frontend_slides_html() -> None:
    plan = """{"slides":[{"title":"收入连续增长","summary":"FY2026收入达到2159亿美元 [Ref 1]","bullets":["数据中心收入是主要增长来源 [Ref 1]"],"chart":{"type":"line","title":"年度营业收入","categories":["FY2024","FY2025","FY2026"],"series":[{"name":"亿美元","values":[609,1305,2159]}],"source_ref":"[Ref 1]"}}]}"""
    item = request().model_copy(update={
        "ppt_skill": "frontend-slides",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=request().sections[0].refs)],
    }, deep=True)
    html = create_html_slides(item).decode("utf-8")

    assert html.startswith("<!doctype html>")
    assert "HTML + JS · 非 PowerPoint 文件" in html
    assert "data-chart=\"0\"" in html
    assert "requestFullscreen" in html and "ArrowRight" in html
    assert "https://" not in html
    assert "[Ref 1]" not in html and "[1]" in html


def test_huawei_skill_uses_different_layouts_for_business_content() -> None:
    sections = [
        DeliverableSection(heading="执行摘要", paragraphs=["库存风险需要同步治理"], bullets=["账实差异需要核验", "长库龄库存需处置", "采购控制需加强"]),
        DeliverableSection(heading="账实差异", paragraphs=["净差异-2.6306万元"], bullets=["盘亏集中于长库龄", "华北仓差异较高", "完成账务核验"]),
        DeliverableSection(heading="高风险物料清单", paragraphs=["SKU008风险最高"], bullets=["SKU015列第二", "SKU002列第三", "SKU004列第四"]),
        DeliverableSection(heading="责任闭环与行动计划", paragraphs=["财务牵头核验"], bullets=["供应链处置", "采购暂停补货", "管理层复核"]),
    ]
    item = request().model_copy(update={"ppt_skill": "guizang-huawei-style-c", "sections": sections}, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    shape_counts = [len(slide.shapes) for slide in list(deck.slides)[1:]]

    assert len(deck.slides) == 5
    assert len(set(shape_counts)) >= 3


def test_huawei_skill_keeps_long_titles_out_of_master_chrome_and_uses_one_font() -> None:
    plan = """{"slides":[{"title":"收入增长强劲但周期波动与资本开支压力仍需同步关注","summary":"2025财年收入与利润明显改善，HBM需求成为增长核心驱动力，但高强度投资可能影响现金流表现","bullets":["数据中心需求持续拉动高带宽存储器收入","传统DRAM与NAND市场仍具有明显周期性"]}]}"""
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c", "include_citations": False,
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=[])],
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    slide = deck.slides[1]
    heading = next(shape for shape in slide.shapes
                   if hasattr(shape, "text_frame") and shape.top == Inches(0.62) and shape.text.strip())
    text_runs = [run for shape in slide.shapes if hasattr(shape, "text_frame")
                 for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text.strip()]

    assert heading.left >= Inches(1.82)
    assert "\n" not in heading.text
    assert len(heading.text) <= 26
    assert all(run.font.name == "Microsoft YaHei" for run in text_runs)


def test_chart_slide_has_one_source_footer() -> None:
    plan = """{"slides":[{"title":"收入连续增长","summary":"收入保持增长 [Ref 1]","bullets":["数据中心收入是主要增长来源 [Ref 1]"],"chart":{"type":"line","title":"年度营业收入","categories":["FY2024","FY2025","FY2026"],"series":[{"name":"亿美元","values":[609,1305,2159]}],"source_ref":"[Ref 1]"}}]}"""
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=request().sections[0].refs)],
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    slide_text = "\n".join(shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text"))

    assert slide_text.count("图表来源：") == 1
    assert "参考：" not in slide_text


def test_pptx_renders_native_chart_and_inline_refs() -> None:
    plan = """{"slides":[{"title":"收入连续五年增长","summary":"FY2026收入达到2159亿美元 [Ref 1]","bullets":["数据中心收入是主要增长来源 [Ref 1]"],"chart":{"type":"line","title":"FY2022-FY2026营业收入","categories":["FY2022","FY2023","FY2024","FY2025","FY2026"],"series":[{"name":"营业收入（亿美元）","values":[269,270,609,1305,2159]}],"source_ref":"[Ref 1]"}}]}"""
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=request().sections[0].refs)],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert any(getattr(shape, "has_chart", False) for slide in deck.slides for shape in slide.shapes)
    assert "[1]" in text
    assert "[Ref 1]" not in text
    assert "经营报告.pdf" in text
    citation_runs = [run for shape in deck.slides[1].shapes if hasattr(shape, "text_frame")
                     for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text == "[1]"]
    assert citation_runs and all(run.font.size.pt <= 10 for run in citation_runs)


def test_word_and_pdf_render_chart_with_source_refs() -> None:
    content = """{"sections":[{"heading":"经营趋势","summary":"收入连续增长 [Ref 1]","bullets":["FY2026达到2159亿美元 [Ref 1]"],"chart":{"type":"bar","title":"年度营业收入","categories":["FY2024","FY2025","FY2026"],"series":[{"name":"亿美元","values":[609,1305,2159]}],"source_ref":"[Ref 1]"}}]}"""
    item = request().model_copy(update={
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[content], bullets=[], refs=request().sections[0].refs)],
    }, deep=True)

    word = Document(BytesIO(create_docx(item)))
    word_text = "\n".join(paragraph.text for paragraph in word.paragraphs)
    pdf = create_pdf(item)
    pdf_text = "\n".join((page.extract_text() or "") for page in PdfReader(BytesIO(pdf)).pages)

    assert len(word.inline_shapes) == 1
    assert "[1]" in word_text and "经营报告.pdf" in word_text
    assert "[Ref 1]" not in word_text
    assert b"/Image" in pdf
    assert "[1]" in pdf_text and "经营报告.pdf" in pdf_text
    assert "[Ref 1]" not in pdf_text


def test_pptx_turns_structured_slide_plan_into_readable_slides() -> None:
    plan = """```json
    {"slides":[
      {"title":"长库龄库存风险集中","summary":"25.1502万元模拟减值需要关注","bullets":["60天以内库存占比39.4%","365天以上库存占比14.0%","华北仓风险最高","优先处置旧版物料"]},
      {"title":"优先处置华北仓风险","summary":"华北仓模拟减值风险最高","bullets":["财务核验盘点差异","供应链制定处置计划"]}
    ]}
    ```"""
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c",
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=[])],
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert len(deck.slides) == 3
    assert "长库龄库存风险集中" in text
    assert "60天以内库存占比39.4%" in text
    assert "25.1502万元" in text
    assert "模拟减值需要关注" in text
    assert "（续）" not in text
    assert "第1页" not in text


def test_pptx_keeps_fourteen_structured_content_slides() -> None:
    slides = [{
        "title": f"第 {index} 个经营判断",
        "summary": f"第 {index} 个判断具有明确的业务影响",
        "bullets": ["量化依据已核对", "影响范围已识别", "下一步动作已明确"],
        "chart": None,
    } for index in range(1, 15)]
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c", "include_citations": False,
        "sections": [DeliverableSection(
            heading="分析结果", paragraphs=[json.dumps({"slides": slides}, ensure_ascii=False)], bullets=[], refs=[]
        )],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))

    assert len(deck.slides) == 15


def test_structured_ppt_keeps_summary_together_and_removes_title_citations() -> None:
    plan = json.dumps({"slides": [{
        "title": "现金转化需重点跟踪 [1]",
        "summary": "经营现金流继续增长；但现金转化率出现回落 [1]",
        "bullets": ["量化依据已核对 [1]", "营运资金压力上升 [1]", "建议按季度跟踪 [1]"],
        "chart": None,
    }]}, ensure_ascii=False)
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c", "include_citations": False,
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[plan], bullets=[], refs=[])],
    }, deep=True)

    deck = Presentation(BytesIO(create_pptx(item)))
    slide_text = "\n".join(shape.text for shape in deck.slides[1].shapes if hasattr(shape, "text"))

    assert len(deck.slides) == 2
    assert "现金转化需重点跟踪" in slide_text
    assert "现金转化需重点跟踪 [1]" not in slide_text
    assert "（续）" not in slide_text


def test_huawei_skill_adds_timeline_priority_and_comparison_layouts() -> None:
    sections = [
        DeliverableSection(heading="执行摘要", paragraphs=["经营优先级已明确"], bullets=["数据已核对", "影响已识别", "动作已安排"]),
        DeliverableSection(heading="年度时间线与里程碑", paragraphs=["一季度完成盘点"], bullets=["二季度完成修复", "三季度进入验收", "四季度完成闭环"]),
        DeliverableSection(heading="区域风险与准备度", paragraphs=["欧盟区域需立即处理"], bullets=["亚太需重点跟踪", "美国需计划改善", "其他区域持续监测"]),
        DeliverableSection(heading="业务结构对比", paragraphs=["数据中心占比领先"], bullets=["游戏业务保持稳定", "汽车业务仍在培育", "专业可视化规模较小"]),
    ]
    item = request().model_copy(update={
        "ppt_skill": "guizang-huawei-style-c", "include_citations": False, "sections": sections,
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    slide_texts = ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
                   for slide in list(deck.slides)[1:]]

    assert "关键节点" in slide_texts[1]
    assert "立即处理" in slide_texts[2] and "持续监测" in slide_texts[2]
    assert "数据中心占比领先" in slide_texts[3]


def test_pptx_recovers_legacy_page_outline_without_showing_page_labels() -> None:
    outline = """第1页｜账实净差异需要核验
账面金额233.28万元
实盘金额230.6494万元
第2页｜超过180天库存占比达33.5%
60天以内库存占比39.4%
365天以上库存占比14.0%"""
    item = request().model_copy(update={
        "sections": [DeliverableSection(heading="分析结果", paragraphs=[outline], bullets=[], refs=[])],
    }, deep=True)
    deck = Presentation(BytesIO(create_pptx(item)))
    text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert len(deck.slides) == 3
    assert "第1页" not in text
    assert "第2页" not in text
    assert "60天以内库存占比39.4%" in text


def test_mermaid_is_valid_flowchart_text() -> None:
    content = create_mermaid(request()).decode("utf-8")
    assert content.startswith("flowchart TD")
    assert "核心结论" in content


def test_mermaid_keeps_llm_generated_source() -> None:
    item = request().model_copy(deep=True)
    item.sections[0].paragraphs = ["flowchart LR\n  A[开始] --> B{是否完成}\n  B -->|Yes| C[输出]"]

    content = create_mermaid(item).decode("utf-8")

    assert content.startswith("flowchart LR")
    assert "B -->|Yes| C" in content


def test_generates_editable_excalidraw_json() -> None:
    item = request().model_copy(deep=True)
    item.sections[0].paragraphs = ["flowchart LR\n  A[开始] --> B{是否完成}\n  B -->|Yes| C[输出]"]

    data = __import__("json").loads(create_excalidraw(item))

    assert data["type"] == "excalidraw"
    assert any(element["type"] == "rectangle" for element in data["elements"])
    assert any(element["type"] == "arrow" for element in data["elements"])
    assert all(element["roughness"] == 2 for element in data["elements"])
